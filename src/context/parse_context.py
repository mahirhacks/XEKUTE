#!/usr/bin/env python3
"""Combine imported assessment material into one bounded Markdown context file."""
from __future__ import annotations

import argparse
import csv
import html
import json
import re
from pathlib import Path

MAX_FILE_BYTES = 8 * 1024 * 1024
MAX_TEXT_CHARS = 200_000


def extract(path: Path) -> tuple[str, str]:
    if path.stat().st_size > MAX_FILE_BYTES:
        return "", "skipped: file exceeds 8 MB"
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            from pypdf import PdfReader  # optional
            return "\n\n".join((page.extract_text() or "") for page in PdfReader(str(path)).pages), "PDF"
        if suffix == ".docx":
            from docx import Document  # optional
            return "\n".join(p.text for p in Document(str(path)).paragraphs), "Word document"
        if suffix in {".xlsx", ".xlsm"}:
            from openpyxl import load_workbook  # optional
            book = load_workbook(path, read_only=True, data_only=True)
            lines = []
            for sheet in book.worksheets:
                lines.append(f"Sheet: {sheet.title}")
                lines.extend(" | ".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))
            return "\n".join(lines), "Excel workbook"
        if suffix == ".pptx":
            from pptx import Presentation  # optional
            deck = Presentation(str(path))
            return "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text")), "PowerPoint presentation"
        if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"}:
            from PIL import Image  # optional
            import pytesseract  # optional
            return pytesseract.image_to_string(Image.open(path)), "image OCR"
        raw = path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".json":
            raw = json.dumps(json.loads(raw), indent=2, ensure_ascii=False)
        elif suffix in {".html", ".htm"}:
            raw = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", raw, flags=re.I | re.S)
            raw = html.unescape(re.sub(r"<[^>]+>", " ", raw))
        elif suffix in {".csv", ".tsv"}:
            delimiter = "\t" if suffix == ".tsv" else ","
            rows = list(csv.reader(raw.splitlines(), delimiter=delimiter))[:2000]
            raw = "\n".join(" | ".join(cell.replace("|", "\\|") for cell in row) for row in rows)
        return raw, suffix.lstrip(".").upper() or "text"
    except Exception as exc:
        return "", f"skipped: {exc}"


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", required=True)
    parser.add_argument("files", nargs="+")
    args = parser.parse_args()
    sections = ["# Penetration Testing Context", "", "Generated from imported source material.", ""]
    parsed = 0
    for raw_path in args.files:
        path = Path(raw_path)
        text, kind = extract(path)
        sections.extend([f"## {path.name}", "", f"_Source: `{path}` · {kind}_", ""])
        if text.strip():
            clean = text.strip()[:MAX_TEXT_CHARS]
            sections.extend(["```text", clean, "```", ""])
            parsed += 1
        else:
            sections.extend([f"> {kind}", ""])
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(sections), encoding="utf-8")
    print(json.dumps({"ok": True, "parsed": parsed, "total": len(args.files), "output": str(output)}))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
